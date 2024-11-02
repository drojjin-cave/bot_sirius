import asyncio
import logging
from aiogram import Bot, Dispatcher, types
from aiogram import F
from aiogram.types import Message
from aiogram.filters import Command

from config import BOT_TOKEN, TOKEN_GROQ

# v2 от 02.11.24

from pptx import Presentation
from llama_index.llms.groq import Groq

llm = Groq(model="llama3-70b-8192", api_key=TOKEN_GROQ)


# Включаем логирование, чтобы не пропустить важные сообщения
logging.basicConfig(level=logging.INFO)
# Объект бота
bot = Bot(token=BOT_TOKEN)
# Диспетчер
dp = Dispatcher()

# Хэндлер на команду /start
@dp.message(Command("start"))
async def cmd_start(message: types.Message):
    await message.answer("Отправьте мне проектную презентацию в формате PPTX файла", reply_markup=types.ReplyKeyboardRemove())


@dp.message(F.document)
async def test(message: Message, bot: Bot):
    #await message.answer("Document received!")
    file_id = message.document.file_id
    file = await bot.get_file(file_id)
    file_path = file.file_path
    file_name = file_id[-10:] + "." + file_path.split('/')[-1].split('.')[-1]
    print(file_name)
    print(file_path)
    await bot.download_file(file_path, file_name)

    if file_name.split(sep='.')[-1].lower() != 'pptx':
      await message.answer('ОШИБКА:\nАнализируются только презентации в формате PPTX!')
      return None

    # Выделяем текст из презентации
    prs = Presentation(file_name)
    slide = prs.slides[0]
    slide_text = ' '
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        slide_text += (shape.text + '\n')
    slide_0_text = slide_text
    #print(slide_0_text)
    slide_goal_text = ' '
    for i, slide in enumerate(prs.slides):
      if i == 0: continue
      for i, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
          if 'цель' in shape.text.lower():
            slide_goal_text += (shape.text + '\n')
            if len(shape.text.split()) < 4:
              if i+1 < len(slide.shapes): slide_goal_text += (slide.shapes[i+1].text + '\n')
            break
    #print(slide_goal_text)


    # Испольуем LLM для выделения сущностей и анализа соответствия
    pres_title = llm.complete('Дан текст с титульной страницы презентации:\n\n' + slide_0_text + '\n Выдели из всего текста только то, что является названием презентации. В ответе выведи ТОЛЬКО название или слово "НЕТ", если название отсутствует.')
    #print(pres_title.text)
    if pres_title.text == 'НЕТ':
      await message.answer('ОШИБКА:\nНа первой странице презентации должно быть текстовое поле с названием проекта.')
      return None
    else:
      await message.answer('Название проекта:\n' + pres_title.text)

    pres_goal = llm.complete('Дан текст со страницы презентации с описанием цели проектной работы:\n\n' + slide_goal_text + '\n Выдели из всего текста только то, что является формулировкой цели. В ответе выведи ТОЛЬКО формулировку цели или слово "НЕТ", если цель отсутствует.')
    #print(pres_goal.text)
    if pres_goal.text == 'НЕТ':
      await message.answer('ОШИБКА:\nВ проекте должна быть описана ОДНА цель, соответствующая названию проекта.')
      return None
    else:
      await message.answer('Цель проекта:\n' + pres_goal.text)

    response = llm.complete("Отвечай ТОЛЬКО на русском языке. Мы проверяем проектную презентацию. Тема презентации " + pres_title.text + ". Цель проекта, указанная в презентации: " + pres_goal.text + ". Предоставь анализ соответствует ли цель проекта указанной теме? ")
    print(response.text)
    #await message.answer("Done!")
    await message.answer('Анализ соответствия темы и цели:\n' + response.text)

# Запуск процесса поллинга новых апдейтов
async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
